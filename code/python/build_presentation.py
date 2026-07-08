"""Build the 'challenges so far' slide deck (16:9) into slides/.

A clean, presentation-ready PowerPoint that reuses the project's own charts to
illustrate each modelling challenge. Run after the charts exist in image/.
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
IMG = ROOT / "image"
SLIDES = ROOT / "slides"
ASSETS = SLIDES / "assets"

INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x89, 0x87, 0x81)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
TEAL = RGBColor(0x1B, 0xAF, 0x7A)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
VIOLET = RGBColor(0x4A, 0x3A, 0xA7)
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
PANEL = RGBColor(0xF2, 0xF1, 0xEC)
LINE = RGBColor(0xE1, 0xE0, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xC3, 0xC2, 0xB7)
FONT = "Calibri"
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         after=6, spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(after); p.line_spacing = spacing
        for txt, size, bold, color in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = FONT
    return tb


def place(slide, path, bx, by, bw, bh, border=True):
    iw, ih = Image.open(path).size
    ar, bar = iw / ih, bw / bh
    if ar > bar:
        w, h = bw, bw / ar
    else:
        h, w = bh, bh * ar
    pic = slide.shapes.add_picture(str(path), Inches(bx + (bw - w) / 2),
                                   Inches(by + (bh - h) / 2), Inches(w), Inches(h))
    if border:
        pic.line.color.rgb = LINE; pic.line.width = Pt(1)
    return pic


def header(slide, section, color, headline, hsize=30):
    rect(slide, 0, 0, SW, 0.16, color)
    text(slide, 0.7, 0.42, 12, 0.4, [[(section.upper(), 13.5, True, color)]])
    text(slide, 0.7, 0.82, 12, 1.1, [[(headline, hsize, True, INK)]])
    text(slide, 0.7, SH - 0.45, 9, 0.35,
         [[("Modelling House Prices in Cambridgeshire", 11, False, MUTED)]])


def bullets(slide, x, y, w, items, color, size=15, after=13):
    paras = [[("▪  ", size - 2, True, color), (b, size, False, INK2)] for b in items]
    text(slide, x, y, w, 4.6, paras, after=after, spacing=1.06)


# ---------- 1. title ----------
s = prs.slides.add_slide(BLANK); bg(s, INK)
iw, ih = Image.open(ASSETS / "kings_college.jpg").size
ph_h = SW * ih / iw
s.shapes.add_picture(str(ASSETS / "kings_college.jpg"), 0, 0, Inches(SW), Inches(ph_h))
rect(s, 0.7, ph_h + 0.55, 1.6, 0.10, BLUE)
text(s, 0.7, ph_h + 0.78, 12.4, 1.2,
     [[("Modelling House Prices in Cambridgeshire", 34, True, WHITE)]])
text(s, 0.7, ph_h + 1.78, 12, 0.6,
     [[("Rapid Report — the challenges so far", 22, False, PALE)]])
text(s, 0.7, SH - 0.85, 12, 0.5,
     [[("Diana · Matt · Davies · Natsani · Xiangjun · Vedanta", 15, False, MUTED)]])

# ---------- 2. the dataset ----------
s = prs.slides.add_slide(BLANK); bg(s, SURFACE)
header(s, "The dataset", BLUE, "What we’re modelling")
tiles = [("59,787", "sold prices", BLUE), ("24", "property types", TEAL),
         ("43", "areas (MSOAs)", ORANGE), ("£350k", "median price", VIOLET)]
tx, tw, gap = 0.7, 2.95, 0.2
for i, (num, lab, col) in enumerate(tiles):
    x = tx + i % 2 * (tw + gap)
    y = 1.95 + i // 2 * 1.5
    rect(s, x, y, tw, 1.32, PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, y + 0.14, tw, 0.7, [[(num, 32, True, col)]], align=PP_ALIGN.CENTER)
    text(s, x, y + 0.86, tw, 0.4, [[(lab, 12.5, False, INK2)]], align=PP_ALIGN.CENTER)
bullets(s, 0.7, 5.15, 6.3, [
    "Sold prices 2018–22 across Cambridge & its districts",
    "Three feature families: property, location, social profile",
    "Rich — but every family brought its own challenges",
], BLUE)
place(s, IMG / "cambridgeshire_price_map.png", 7.35, 1.5, 5.6, 5.7)

# ---------- 3. challenge 1: data quality ----------
s = prs.slides.add_slide(BLANK); bg(s, SURFACE)
header(s, "Challenge 1 · the data fights back", ORANGE, "Time, duplicates & missing values")
bullets(s, 0.7, 1.95, 12.0, [
    "Sold dates imputed to 31 March (year-end) — not real timing, and not a time series",
    "59% of rows are exact duplicates — must de-duplicate before modelling",
    "Bedrooms / bathrooms coded 0 = missing (12k baths fixed, 16k receptions to handle)",
], ORANGE, after=8)
place(s, IMG / "cambridgeshire_sales_timeline.png", 0.7, 3.55, 12.0, 3.55)

# ---------- 4. challenge 2: dimensionality ----------
s = prs.slides.add_slide(BLANK); bg(s, SURFACE)
header(s, "Challenge 2 · too many variables", TEAL, "High-dimensional — and highly redundant")
bullets(s, 0.7, 2.05, 6.1, [
    "24 types + ~25 census variables → high-dimensional",
    "Census blocks (age, commute, quals) each sum to 100% → collinear",
    "Selection: intuitive vs optimised; Stepwise vs RFE",
    "Baseline linear vs complex non-linear models",
], TEAL, after=16)
place(s, IMG / "cambridgeshire_correlation_heatmap.png", 7.0, 1.2, 6.0, 5.9)

# ---------- 5. challenge 3: what to keep & how ----------
s = prs.slides.add_slide(BLANK); bg(s, SURFACE)
header(s, "Challenge 3 · what to keep & how", VIOLET, "Location already carries the social signal")
bullets(s, 0.7, 1.95, 12.0, [
    "Strong social→price links collapse ~64% once distance to Cambridge & London is known",
    "So we model property + location only — no social variables, almost no signal lost",
    "Target is skewed & bounded → model log(price)",
], VIOLET, after=8)
place(s, IMG / "cambridgeshire_social_vs_location.png", 1.5, 3.5, 10.3, 3.6)

out = SLIDES / "rapid_report_challenges.pptx"
prs.save(str(out))
print(f"Saved {out.relative_to(ROOT)}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
